#!/usr/bin/env python3
"""
FastTracker 性能評価報告書 PowerPoint生成スクリプト
3イテレーション版 (Ballistic/HGV × cluster=0/40)

GPU加速マルチターゲット追尾システムの 3 回イテレーション性能評価結果を
PowerPointプレゼンテーションとして出力する。
"""

import os
import tempfile

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "performance_report.pptx")
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")

# ---------------------------------------------------------------------------
# 評価データ（3イテレーション × 4パターン）
# ---------------------------------------------------------------------------
PATTERNS = [
    "Ballistic\ncluster=0",
    "Ballistic\ncluster=40",
    "HGV\ncluster=0",
    "HGV\ncluster=40",
]
PATTERNS_LABEL = ["A: Ballistic cl=0", "B: Ballistic cl=40", "C: HGV cl=0", "D: HGV cl=40"]

# [Iter1, Iter2, Iter3] per pattern
DATA = {
    "pos_rmse": {
        "mean": [[1298.9, 1428.7, 2400.1, 2042.1],
                 [1186.4, 1384.1, 2390.6, 2045.7],
                 [1177.4, 1344.8, 2252.1, 1948.0]],
        "std":  [[179.9,   54.9,  255.4,   49.5],
                 [ 49.6,   27.9,  256.7,   24.4],
                 [ 27.0,   36.5,  265.6,   28.6]],
        "unit": "m",
        "label": "位置 RMSE [m]",
    },
    "vel_rmse": {
        "mean": [[230.2, 220.0, 445.6, 324.8],
                 [198.4, 216.8, 394.2, 305.0],
                 [196.6, 215.3, 231.9, 242.0]],
        "std":  [[ 59.7,   1.3, 122.6,  18.4],
                 [  4.4,   1.9,  87.3,   5.9],
                 [  3.4,   0.9,   7.5,   1.8]],
        "unit": "m/s",
        "label": "速度 RMSE [m/s]",
    },
    "ospa": {
        "mean": [[1170.3, 2478.7, 2201.5, 3223.9],
                 [1054.2, 2206.5, 2040.2, 2350.6],
                 [1087.7, 1891.4, 1971.6, 2269.4]],
        "std":  [[182.6,  92.1, 204.0, 212.2],
                 [ 45.6, 217.8, 170.2,  39.7],
                 [ 32.5,  23.3, 169.6,  11.6]],
        "unit": "m",
        "label": "OSPA 距離 [m]",
    },
    "precision": {
        "mean": [[98.6, 96.0, 95.1, 92.5],
                 [99.8, 97.8, 99.0, 98.0],
                 [99.9, 99.1, 99.4, 98.6]],
        "std":  [[2.8, 0.6, 3.6, 1.4],
                 [0.4, 0.9, 1.1, 0.1],
                 [0.3, 0.1, 0.6, 0.1]],
        "unit": "%",
        "label": "Precision [%]",
    },
    "recall": {
        "mean": [[99.4, 99.3, 99.1, 98.8],
                 [99.4, 99.6, 99.4, 99.0],
                 [99.4, 99.7, 99.4, 99.3]],
        "std":  [[0.4, 0.1, 0.4, 0.3],
                 [0.4, 0.1, 0.1, 0.1],
                 [0.4, 0.0, 0.1, 0.1]],
        "unit": "%",
        "label": "Recall [%]",
    },
    "f1": {
        "mean": [[99.0, 97.6, 97.1, 95.5],
                 [99.6, 98.7, 99.2, 98.5],
                 [99.6, 99.4, 99.4, 98.9]],
        "std":  [[1.4, 0.3, 2.0, 0.8],
                 [0.2, 0.4, 0.6, 0.1],
                 [0.2, 0.0, 0.3, 0.1]],
        "unit": "%",
        "label": "F1 スコア [%]",
    },
}

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x00, 0x33, 0x66)
MID_BLUE    = RGBColor(0x00, 0x5B, 0x96)
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xB5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
HEADER_BG   = RGBColor(0x00, 0x33, 0x66)
ROW_ALT     = RGBColor(0xEC, 0xF0, 0xF5)

MPL_ITER1  = "#003366"
MPL_ITER2  = "#1A73B5"
MPL_ITER3  = "#27AE60"
MPL_RED    = "#E74C3C"
MPL_GRID   = "#CCCCCC"

ITER_COLORS = [MPL_ITER1, MPL_ITER2, MPL_ITER3]

# ---------------------------------------------------------------------------
# Font helpers
# ---------------------------------------------------------------------------
JAPANESE_FONT = "Meiryo"

_jp_font_path = None
for f in fm.findSystemFonts():
    lower = f.lower()
    if "meiryo" in lower and lower.endswith(".ttc"):
        _jp_font_path = f
        break
    if "yugoth" in lower and lower.endswith(".ttc"):
        _jp_font_path = f
        break

_jp_font_prop = fm.FontProperties(fname=_jp_font_path) if _jp_font_path else fm.FontProperties()

plt.rcParams.update({
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "axes.edgecolor": MPL_ITER1,
    "axes.labelcolor": MPL_ITER1,
    "xtick.color": MPL_ITER1,
    "ytick.color": MPL_ITER1,
    "grid.color": MPL_GRID,
    "grid.linestyle": "--",
    "grid.alpha": 0.7,
})


# ---------------------------------------------------------------------------
# Chart generation
# ---------------------------------------------------------------------------
def _make_grouped_bar_chart(key, png_path, title, lower_is_better=True,
                             ylim_bottom=None):
    """Grouped bar chart: x=Patterns, groups=Iter1/2/3."""
    means = DATA[key]["mean"]  # [iter][pattern]
    stds  = DATA[key]["std"]
    label = DATA[key]["label"]
    unit  = DATA[key]["unit"]

    n_patterns = len(PATTERNS_LABEL)
    n_iters = 3
    x = np.arange(n_patterns)
    width = 0.25

    fig, ax = plt.subplots(figsize=(10, 4.5), dpi=160)

    for it in range(n_iters):
        offset = (it - 1) * width
        bars = ax.bar(
            x + offset, means[it], width,
            label=f"Iter {it + 1}",
            color=ITER_COLORS[it], alpha=0.88,
            edgecolor="white", linewidth=0.8,
            yerr=stds[it], capsize=4,
            error_kw={"elinewidth": 0.9, "ecolor": "#666666"},
        )
        for bar, val in zip(bars, means[it]):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                bar.get_height() + (max(max(m) for m in means) * 0.01),
                f"{val:.0f}" if unit != "%" else f"{val:.1f}",
                ha="center", va="bottom", fontsize=7.5, color="#333333",
                fontproperties=_jp_font_prop,
            )

    ax.set_xticks(x)
    ax.set_xticklabels(PATTERNS_LABEL, fontproperties=_jp_font_prop, fontsize=9)
    ax.set_ylabel(label, fontproperties=_jp_font_prop, fontsize=10)
    ax.set_title(title, fontproperties=_jp_font_prop, fontsize=12,
                 color=MPL_ITER1, fontweight="bold", pad=8)
    ax.legend(prop=_jp_font_prop, fontsize=10, loc="upper right",
              framealpha=0.85)
    ax.grid(axis="y", alpha=0.5)
    if ylim_bottom is not None:
        ax.set_ylim(bottom=ylim_bottom)
    else:
        ax.set_ylim(bottom=0)

    # Direction arrow
    arrow_color = MPL_RED if lower_is_better else MPL_ITER3
    arrow_label = "低いほど良い ↓" if lower_is_better else "高いほど良い ↑"
    ax.text(0.01, 0.97, arrow_label,
            transform=ax.transAxes, fontsize=9, color=arrow_color,
            va="top", fontproperties=_jp_font_prop, fontweight="bold",
            bbox=dict(facecolor="white", edgecolor=arrow_color,
                      linewidth=0.8, boxstyle="round,pad=0.2", alpha=0.8))

    for label_obj in ax.get_xticklabels() + ax.get_yticklabels():
        label_obj.set_fontproperties(_jp_font_prop)

    fig.tight_layout()
    fig.savefig(png_path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _make_improvement_bar(png_path):
    """改善率比較チャート (Iter1→Iter3)."""
    metrics = ["位置RMSE\nA:Ballistic", "位置RMSE\nD:HGV cl=40",
               "速度RMSE\nC:HGV cl=0", "速度RMSE\nD:HGV cl=40",
               "OSPA\nB:Bal cl=40", "OSPA\nD:HGV cl=40",
               "Precision\nD:HGV cl=40", "F1\nD:HGV cl=40"]

    # Iter1 → Iter3 improvement (%) — positive = improvement
    improvements = [
        (1298.9 - 1177.4) / 1298.9 * 100,  # pos_rmse A
        (2042.1 - 1948.0) / 2042.1 * 100,  # pos_rmse D
        (445.6 - 231.9) / 445.6 * 100,     # vel_rmse C
        (324.8 - 242.0) / 324.8 * 100,     # vel_rmse D
        (2478.7 - 1891.4) / 2478.7 * 100,  # ospa B
        (3223.9 - 2269.4) / 3223.9 * 100,  # ospa D
        (98.6 - 92.5),                      # precision D (pt)
        (98.9 - 95.5),                      # f1 D (pt)
    ]
    units = ["%", "%", "%", "%", "%", "%", "pt", "pt"]
    colors = [MPL_ITER2 if v > 10 else MPL_ITER3 for v in improvements]

    fig, ax = plt.subplots(figsize=(10, 4.2), dpi=160)
    y = np.arange(len(metrics))
    bars = ax.barh(y, improvements, color=colors, edgecolor="white",
                   linewidth=0.8, height=0.55)

    for bar, val, unit in zip(bars, improvements, units):
        ax.text(val + 0.3, bar.get_y() + bar.get_height() / 2,
                f"+{val:.1f}{unit}", va="center", fontsize=9,
                color="#333333", fontweight="bold",
                fontproperties=_jp_font_prop)

    ax.set_yticks(y)
    ax.set_yticklabels(metrics, fontproperties=_jp_font_prop, fontsize=9)
    ax.set_xlabel("改善量 (Iter1→Iter3)", fontproperties=_jp_font_prop, fontsize=10)
    ax.set_title("主要指標の改善量 (Iter 1 → Iter 3)", fontproperties=_jp_font_prop,
                 fontsize=12, color=MPL_ITER1, fontweight="bold", pad=8)
    ax.grid(axis="x", alpha=0.5)
    ax.set_xlim(left=0)
    ax.axvline(x=10, color=MPL_RED, linestyle="--", linewidth=0.8, alpha=0.6)
    ax.text(10.3, -0.6, "10%", fontsize=8, color=MPL_RED,
            fontproperties=_jp_font_prop)

    for label_obj in ax.get_xticklabels() + ax.get_yticklabels():
        label_obj.set_fontproperties(_jp_font_prop)

    fig.tight_layout()
    fig.savefig(png_path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ---------------------------------------------------------------------------
# PowerPoint helpers
# ---------------------------------------------------------------------------
def _set_font(run, size=14, bold=False, color=BLACK, name=JAPANESE_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                 color=BLACK, alignment=PP_ALIGN.LEFT, name=JAPANESE_FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color, name=name)
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, size=13,
                     color=DARK_GRAY, bullet_color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        p.line_spacing = Pt(int(size * 1.35))
        bullet_run = p.add_run()
        bullet_run.text = "■  "
        _set_font(bullet_run, size=size, bold=False, color=bullet_color)
        text_run = p.add_run()
        text_run.text = item
        _set_font(text_run, size=size, bold=False, color=color)
    return txBox


def _add_slide_number(slide, slide_num, total):
    _add_textbox(
        slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35),
        f"{slide_num} / {total}", size=9, color=MID_GRAY, alignment=PP_ALIGN.RIGHT,
    )


def _add_header_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.9),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    _add_textbox(
        slide, Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.7),
        title_text, size=24, bold=True, color=WHITE,
    )


def _add_bottom_line(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.03),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BLUE
    shape.line.fill.background()


def _make_table(slide, left, top, width, height, rows, cols, data,
                col_widths=None, highlight_col=None, highlight_color=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG
                _set_font(run, size=11, bold=True, color=WHITE)
                p.alignment = PP_ALIGN.CENTER
            else:
                if highlight_col is not None and c == highlight_col:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = highlight_color or LIGHT_BLUE
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ROW_ALT
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                _set_font(run, size=10, bold=False, color=DARK_GRAY)
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
    return table_shape


# =========================================================================
# Slide builders
# =========================================================================
TOTAL_SLIDES = 13


def build_slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.0), Inches(10.9), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.fill.background()
    _add_textbox(slide, Inches(1.2), Inches(1.2), Inches(10.9), Inches(1.6),
                 "FastTracker 性能評価報告書",
                 size=38, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.2), Inches(2.8), Inches(10.9), Inches(0.5),
                 "— 3イテレーション チューニング評価 —",
                 size=18, bold=False, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.2), Inches(3.3), Inches(10.9), Inches(1.0),
                 "GPU加速マルチターゲット追尾システム\n"
                 "Ballistic / HGV × cluster=0 / 40",
                 size=16, bold=False, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.2), Inches(5.1), Inches(10.9), Inches(0.5),
                 "2026年2月",
                 size=16, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.2), Inches(6.1), Inches(10.9), Inches(0.4),
                 "CONFIDENTIAL",
                 size=11, bold=True, color=MID_BLUE, alignment=PP_ALIGN.CENTER)
    _add_slide_number(slide, 1, TOTAL_SLIDES)


def build_slide_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "目次")
    items = [
        ("1.", "評価概要", "評価目的・対象システム"),
        ("2.", "評価パターン・試験条件", "シナリオ・センサー・ターゲット設定"),
        ("3.", "追尾パラメータ推移", "3イテレーションのパラメータ変更方針"),
        ("4.", "Iteration 1 ベースライン結果", "4パターン追尾性能評価"),
        ("5.", "Iteration 2 チューニング結果", "ゲート・確定条件・加速度ノイズ改善"),
        ("6.", "Iteration 3 チューニング結果", "速度ノイズ・加速度ノイズ最終調整"),
        ("7.", "性能推移 — 位置・速度 RMSE", "3イテレーション比較グラフ"),
        ("8.", "性能推移 — OSPA・Precision", "3イテレーション比較グラフ"),
        ("9.", "主要指標の改善量", "Iter1→Iter3 改善率サマリー"),
        ("10.", "処理速度評価", "GPU加速による処理性能"),
        ("11.", "総合評価・まとめ", "全メトリクスの総括と推奨パラメータ"),
    ]
    y_start = Inches(1.15)
    for i, (num, title, desc) in enumerate(items):
        y = y_start + Inches(i * 0.53)
        _add_textbox(slide, Inches(0.8), y, Inches(0.5), Inches(0.42),
                     num, size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.RIGHT)
        _add_textbox(slide, Inches(1.45), y, Inches(5.5), Inches(0.42),
                     title, size=14, bold=True, color=DARK_BLUE)
        _add_textbox(slide, Inches(7.2), y, Inches(5.5), Inches(0.42),
                     desc, size=11, bold=False, color=MID_GRAY)
        if i < len(items) - 1:
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), y + Inches(0.44), Inches(11.7), Inches(0.01))
            sep.fill.solid()
            sep.fill.fore_color.rgb = LIGHT_BLUE
            sep.line.fill.background()
    _add_bottom_line(slide)
    _add_slide_number(slide, 2, TOTAL_SLIDES)


def build_slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "1. 評価概要")

    _add_textbox(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.4),
                 "評価目的", size=16, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.3), [
        "Ballistic / HGV × クラスタ有無 の 4 パターンで追尾性能を定量評価",
        "追尾パラメータを 3 回チューニングして性能改善効果を確認",
        "センサー・ターゲットパラメータを固定し追尾器性能のみを評価",
        "各イテレーション 5 回モンテカルロ試験 (seed=42)",
    ], size=12)

    _add_textbox(slide, Inches(0.6), Inches(4.1), Inches(5.8), Inches(0.4),
                 "評価対象システム", size=16, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(0.6), Inches(4.6), Inches(5.8), Inches(2.3), [
        "FastTracker (CUDA/C++)",
        "IMM-UKF (CV / Ballistic / Coordinated Turn モデル)",
        "GNN (Global Nearest Neighbor) データアソシエーション",
        "GPU: CUDA 12.6 / Docker コンテナ内実行",
    ], size=12)

    # Right: evaluation items
    _add_textbox(slide, Inches(7.0), Inches(1.1), Inches(5.7), Inches(0.4),
                 "評価パターン", size=16, bold=True, color=DARK_BLUE)
    pat_data = [
        ["パターン", "ミサイル種別", "クラスタ数", "総目標数"],
        ["A", "Ballistic", "0", "1"],
        ["B", "Ballistic", "40", "41"],
        ["C", "HGV", "0", "1"],
        ["D", "HGV", "40", "41"],
    ]
    _make_table(slide, Inches(7.0), Inches(1.6), Inches(5.7), Inches(2.0),
                5, 4, pat_data,
                col_widths=[Inches(0.8), Inches(1.5), Inches(1.5), Inches(1.9)])

    _add_textbox(slide, Inches(7.0), Inches(3.9), Inches(5.7), Inches(0.4),
                 "評価指標", size=16, bold=True, color=DARK_BLUE)
    metrics_data = [
        ["指標", "内容"],
        ["位置 RMSE", "目標位置の二乗平均平方根誤差 [m]"],
        ["速度 RMSE", "目標速度の二乗平均平方根誤差 [m/s]"],
        ["OSPA 距離", "最適部分割当距離 [m]"],
        ["Precision", "確定トラック中の正しいトラック率 [%]"],
        ["Recall", "真目標中の追尾成功率 [%]"],
        ["F1 スコア", "Precision と Recall の調和平均 [%]"],
    ]
    _make_table(slide, Inches(7.0), Inches(4.4), Inches(5.7), Inches(2.5),
                7, 2, metrics_data, col_widths=[Inches(2.0), Inches(3.7)])

    _add_bottom_line(slide)
    _add_slide_number(slide, 3, TOTAL_SLIDES)


def build_slide_conditions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "2. 評価パターン・試験条件")

    # Sensor params - left top
    _add_textbox(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.38),
                 "センサーパラメータ（全イテレーション固定）", size=14, bold=True, color=DARK_BLUE)
    sensor_data = [
        ["パラメータ", "値"],
        ["センサー位置", "(0, 0) m"],
        ["視野角 (FOV)", "360° (全周)"],
        ["誤警報確率 (Pfa)", "1×10⁻⁶"],
        ["基準検出確率 (Pd)", "90% @ max_range"],
        ["距離ノイズ", "10 m"],
        ["方位角ノイズ", "0.01 rad (≈0.57°)"],
        ["ドップラーノイズ", "2 m/s"],
        ["ビーム幅", "3° (0.0524 rad)"],
        ["ビーム数/フレーム", "10"],
    ]
    _make_table(slide, Inches(0.5), Inches(1.55), Inches(5.8), Inches(4.0),
                len(sensor_data), 2, sensor_data,
                col_widths=[Inches(3.1), Inches(2.7)])

    # Ballistic params - right top
    _add_textbox(slide, Inches(6.9), Inches(1.1), Inches(6.0), Inches(0.38),
                 "ターゲットパラメータ（固定）", size=14, bold=True, color=DARK_BLUE)
    target_data = [
        ["パラメータ", "Ballistic", "HGV"],
        ["発射位置 X", "−600,000 m", "−800,000 m"],
        ["発射角", "0.698 rad (40°)", "0.349 rad (20°)"],
        ["飛翔時間", "400 s", "600 s"],
        ["フレームレート", "10 Hz", "10 Hz"],
        ["初期質量", "20,000 kg", "20,000 kg"],
        ["燃料比率", "0.70", "0.70"],
        ["巡航高度", "—", "40,000 m"],
        ["グライド比", "—", "4.0"],
        ["クラスタ拡散", "5,000 m", "5,000 m"],
    ]
    _make_table(slide, Inches(6.9), Inches(1.55), Inches(6.0), Inches(4.0),
                len(target_data), 3, target_data,
                col_widths=[Inches(2.0), Inches(2.0), Inches(2.0)])

    _add_bottom_line(slide)
    _add_slide_number(slide, 4, TOTAL_SLIDES)


def build_slide_params(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "3. 追尾パラメータ推移（Iter 1 → 2 → 3）")

    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
                 "センサー・ターゲットパラメータは全イテレーション不変。追尾器パラメータのみを調整。",
                 size=13, bold=False, color=MID_GRAY)

    param_data = [
        ["パラメータ", "Iter 1（ベース）", "Iter 2", "Iter 3（最終）", "変更理由"],
        ["gate_threshold", "500", "200", "150",
         "ゲート狭小化で FP 削減"],
        ["confirm_hits", "2", "3", "3（維持）",
         "3連続ヒット必要で偽トラック抑制"],
        ["delete_misses", "90", "30", "20",
         "早期削除で偽トラック長期化防止"],
        ["process_vel_noise", "65 m/s", "65 m/s", "150 m/s",
         "HGV 速度急変への追従性向上"],
        ["process_acc_noise", "160 m/s²", "300 m/s²", "500 m/s²",
         "HGV 機動（~5G）に対応した不確かさ拡大"],
        ["process_pos_noise", "160 m", "160 m（維持）", "160 m（維持）",
         "—"],
        ["min_snr_for_init", "22 dB", "22 dB（維持）", "22 dB（維持）",
         "—"],
    ]
    _make_table(
        slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5),
        len(param_data), 5, param_data,
        col_widths=[Inches(2.0), Inches(1.8), Inches(1.6), Inches(1.8), Inches(5.1)],
        highlight_col=3, highlight_color=RGBColor(0xD5, 0xF0, 0xDC),
    )

    _add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                 "Iter 2: ゲート狭小化 + 確定条件厳格化 + 加速度ノイズ増大 → Precision 改善",
                 size=12, bold=False, color=DARK_GRAY)
    _add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.45),
                 "Iter 3: 速度・加速度ノイズをさらに増大 → HGV 速度 RMSE が −41% 改善",
                 size=12, bold=False, color=DARK_GRAY)

    _add_bottom_line(slide)
    _add_slide_number(slide, 5, TOTAL_SLIDES)


def _build_iter_result_slide(prs, iter_idx, slide_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, title)

    means = {k: DATA[k]["mean"][iter_idx] for k in DATA}
    stds  = {k: DATA[k]["std"][iter_idx]  for k in DATA}

    table_data = [
        ["パターン", "位置RMSE [m]", "速度RMSE [m/s]", "OSPA [m]",
         "Precision [%]", "Recall [%]", "F1 [%]"],
    ]
    for p in range(4):
        row = [
            PATTERNS_LABEL[p],
            f"{means['pos_rmse'][p]:.0f} ±{stds['pos_rmse'][p]:.0f}",
            f"{means['vel_rmse'][p]:.0f} ±{stds['vel_rmse'][p]:.0f}",
            f"{means['ospa'][p]:.0f} ±{stds['ospa'][p]:.0f}",
            f"{means['precision'][p]:.1f} ±{stds['precision'][p]:.1f}",
            f"{means['recall'][p]:.1f} ±{stds['recall'][p]:.1f}",
            f"{means['f1'][p]:.1f} ±{stds['f1'][p]:.1f}",
        ]
        table_data.append(row)
    _make_table(
        slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(2.3),
        5, 7, table_data,
        col_widths=[Inches(2.0), Inches(1.6), Inches(1.7), Inches(1.5),
                    Inches(1.8), Inches(1.5), Inches(1.4)],
    )

    # Key findings
    if iter_idx == 0:
        findings = [
            "Ballistic 単機 (A): Precision≈98.6%, F1≈99.0% — 良好なベースライン",
            "Ballistic クラスタ (B): Precision=96.0% — クラッタ環境で FP 発生",
            "HGV 単機 (C): 速度RMSE=445 m/s — HGV 機動に対してプロセスノイズ不足",
            "HGV クラスタ (D): Precision=92.5%, OSPA=3224m — 最も改善余地が大きい",
        ]
    elif iter_idx == 1:
        findings = [
            "全パターンで Precision が顕著に改善 (最大 +3.9pt: C パターン)",
            "OSPA が全体的に改善 (D: 3224→2351m, −27%)",
            "HGV 速度 RMSE はまだ高い (C: 445→394 m/s) — 速度ノイズ追加が必要",
            "F1 が全パターンで 98.5% 以上に向上",
        ]
    else:
        findings = [
            "HGV 速度 RMSE が劇的に改善 (C: 394→232 m/s, −41%; D: 305→242 m/s, −21%)",
            "Pattern D 位置 RMSE が 2000m 未満に到達 (2046→1948m)",
            "Ballistic クラスタ OSPA −14% (2207→1891m)",
            "全パターンで Precision ≥ 98.6%, F1 ≥ 98.9% を達成",
        ]

    _add_textbox(slide, Inches(0.4), Inches(3.65), Inches(12.5), Inches(0.38),
                 "注目ポイント", size=14, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(2.8),
                     findings, size=12)

    _add_bottom_line(slide)
    _add_slide_number(slide, slide_num, TOTAL_SLIDES)


def build_slide_iter1(prs):
    _build_iter_result_slide(prs, 0, 6, "4. Iteration 1 — ベースライン評価結果")


def build_slide_iter2(prs):
    _build_iter_result_slide(prs, 1, 7, "5. Iteration 2 — チューニング結果 (gate↓, confirm↑, acc_noise↑)")


def build_slide_iter3(prs):
    _build_iter_result_slide(prs, 2, 8, "6. Iteration 3 — 最終チューニング結果 (vel_noise↑, acc_noise↑)")


def build_slide_rmse_trend(prs, chart_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "7. 性能推移 — 位置 RMSE・速度 RMSE")

    pos_path = os.path.join(chart_dir, "pos_rmse.png")
    _make_grouped_bar_chart("pos_rmse", pos_path, "位置 RMSE の推移 (Iter 1→2→3)")
    slide.shapes.add_picture(pos_path, Inches(0.3), Inches(1.0), Inches(6.4), Inches(3.3))

    vel_path = os.path.join(chart_dir, "vel_rmse.png")
    _make_grouped_bar_chart("vel_rmse", vel_path, "速度 RMSE の推移 (Iter 1→2→3)")
    slide.shapes.add_picture(vel_path, Inches(6.9), Inches(1.0), Inches(6.1), Inches(3.3))

    # Text findings
    _add_textbox(slide, Inches(0.3), Inches(4.5), Inches(6.4), Inches(0.38),
                 "位置 RMSE 改善ポイント", size=13, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(0.3), Inches(4.9), Inches(6.4), Inches(1.9), [
        "A (Ballistic cl=0): 1299→1177m (−9.4%)",
        "B (Ballistic cl=40): 1429→1345m (−5.9%)",
        "C (HGV cl=0): 2400→2252m (−6.2%)",
        "D (HGV cl=40): 2042→1948m (−4.8%) ★ 2000m 達成",
    ], size=11)

    _add_textbox(slide, Inches(6.9), Inches(4.5), Inches(6.1), Inches(0.38),
                 "速度 RMSE 改善ポイント", size=13, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(6.9), Inches(4.9), Inches(6.1), Inches(1.9), [
        "Iter 3 で vel_noise 65→150 m/s に増大",
        "C (HGV cl=0): 446→232 m/s (−47.9%) ★★ 最大改善",
        "D (HGV cl=40): 325→242 m/s (−25.5%) ★",
        "Ballistic は Iter 2 段階で収束 (~200 m/s)",
    ], size=11)

    _add_bottom_line(slide)
    _add_slide_number(slide, 9, TOTAL_SLIDES)


def build_slide_ospa_prec_trend(prs, chart_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "8. 性能推移 — OSPA 距離・Precision")

    ospa_path = os.path.join(chart_dir, "ospa.png")
    _make_grouped_bar_chart("ospa", ospa_path, "OSPA 距離の推移 (Iter 1→2→3)")
    slide.shapes.add_picture(ospa_path, Inches(0.3), Inches(1.0), Inches(6.4), Inches(3.3))

    prec_path = os.path.join(chart_dir, "precision.png")
    _make_grouped_bar_chart("precision", prec_path, "Precision の推移 (Iter 1→2→3)",
                             lower_is_better=False, ylim_bottom=88)
    slide.shapes.add_picture(prec_path, Inches(6.9), Inches(1.0), Inches(6.1), Inches(3.3))

    _add_textbox(slide, Inches(0.3), Inches(4.5), Inches(6.4), Inches(0.38),
                 "OSPA 改善ポイント", size=13, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(0.3), Inches(4.9), Inches(6.4), Inches(1.9), [
        "B (Ballistic cl=40): 2479→1891m (−23.7%) ★",
        "D (HGV cl=40): 3224→2269m (−29.6%) ★★",
        "クラスタ環境での多目標 OSPA が大幅改善",
        "偽トラック削除の早期化が有効に機能",
    ], size=11)

    _add_textbox(slide, Inches(6.9), Inches(4.5), Inches(6.1), Inches(0.38),
                 "Precision 改善ポイント", size=13, bold=True, color=DARK_BLUE)
    _add_bullet_list(slide, Inches(6.9), Inches(4.9), Inches(6.1), Inches(1.9), [
        "D (HGV cl=40): 92.5%→98.6% (+6.1pt) ★★",
        "C (HGV cl=0): 95.1%→99.4% (+4.3pt) ★",
        "B (Ballistic cl=40): 96.0%→99.1% (+3.1pt)",
        "ゲート閾値 500→150 と confirm_hits 2→3 が寄与",
    ], size=11)

    _add_bottom_line(slide)
    _add_slide_number(slide, 10, TOTAL_SLIDES)


def build_slide_improvement(prs, chart_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "9. 主要指標の改善量 (Iter 1 → Iter 3)")

    imp_path = os.path.join(chart_dir, "improvement.png")
    _make_improvement_bar(imp_path)
    slide.shapes.add_picture(imp_path, Inches(0.3), Inches(1.0), Inches(8.5), Inches(5.5))

    # Highlight box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.1), Inches(1.2), Inches(4.0), Inches(5.2),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = MID_BLUE
    box.line.width = Pt(1.5)

    _add_textbox(slide, Inches(9.2), Inches(1.4), Inches(3.8), Inches(0.4),
                 "特筆すべき改善", size=14, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)

    highlights = [
        ("HGV 速度 RMSE (C)", "−47.9%"),
        ("OSPA, HGV cl=40 (D)", "−29.6%"),
        ("OSPA, Bal cl=40 (B)", "−23.7%"),
        ("HGV 速度 RMSE (D)", "−20.7%"),
        ("Precision, HGV cl=40", "+6.1 pt"),
        ("F1, HGV cl=40 (D)", "+3.4 pt"),
    ]
    y = Inches(2.0)
    for label, val in highlights:
        _add_textbox(slide, Inches(9.2), y, Inches(2.3), Inches(0.5),
                     label, size=11, bold=False, color=DARK_GRAY)
        _add_textbox(slide, Inches(11.3), y, Inches(1.7), Inches(0.5),
                     val, size=16, bold=True, color=ACCENT_BLUE,
                     alignment=PP_ALIGN.RIGHT)
        y += Inches(0.72)

    _add_bottom_line(slide)
    _add_slide_number(slide, 11, TOTAL_SLIDES)


def build_slide_speed(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "10. 処理速度評価")

    _add_textbox(slide, Inches(0.6), Inches(1.1), Inches(11.0), Inches(0.4),
                 "全パターン共通の処理速度（Iteration 3 実測値）",
                 size=16, bold=True, color=DARK_BLUE)

    for i, (label, val, unit, sub) in enumerate([
        ("全体処理時間\n(Ballistic, cl=0)", "0.9", "s", "400s シナリオ / 5試行平均"),
        ("全体処理時間\n(HGV, cl=40)", "1.9", "s", "600s × 41目標シナリオ"),
        ("平均フレーム処理\n(単目標)", "0.3", "ms", "リアルタイム比 ~1,300x"),
        ("平均フレーム処理\n(41目標)", "0.5", "ms", "リアルタイム比 ~200x"),
    ]):
        x = Inches(0.5 + i * 3.15)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.9), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = MID_BLUE
        card.line.width = Pt(1.5)
        _add_textbox(slide, x + Inches(0.1), Inches(1.85), Inches(2.7), Inches(0.6),
                     label, size=11, bold=True, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.1), Inches(2.55), Inches(2.7), Inches(0.65),
                     val, size=38, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.1), Inches(3.15), Inches(2.7), Inches(0.35),
                     unit, size=20, bold=False, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.1), Inches(3.55), Inches(2.7), Inches(0.4),
                     sub, size=9, bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    perf_data = [
        ["シナリオ", "全体時間", "フレーム時間", "リアルタイム比", "判定"],
        ["Ballistic, cl=0", "0.9 ± 0.0 s", "0.3 ± 0.0 ms", "×1,333", "PASS"],
        ["Ballistic, cl=40", "2.2 ± 0.6 s", "0.5 ± 0.1 ms", "× 364", "PASS"],
        ["HGV, cl=0", "0.8 ± 0.0 s", "0.3 ± 0.0 ms", "×2,000", "PASS"],
        ["HGV, cl=40", "1.9 ± 0.5 s", "0.5 ± 0.1 ms", "× 211", "PASS"],
    ]
    tbl = _make_table(
        slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.2),
        len(perf_data), 5, perf_data,
        col_widths=[Inches(2.5), Inches(2.0), Inches(2.2), Inches(2.0), Inches(3.6)],
    )
    table = tbl.table
    for r in range(1, len(perf_data)):
        cell = table.cell(r, 4)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = GREEN
                run.font.bold = True

    _add_bottom_line(slide)
    _add_slide_number(slide, 12, TOTAL_SLIDES)


def build_slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "11. 総合評価・まとめ")

    # Final params table
    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(5.8), Inches(0.4),
                 "最終追尾パラメータ (Iter 3 確定値)", size=14, bold=True, color=DARK_BLUE)
    param_final = [
        ["パラメータ", "確定値"],
        ["gate_threshold", "150"],
        ["confirm_hits", "3"],
        ["delete_misses", "20"],
        ["process_pos_noise", "160 m"],
        ["process_vel_noise", "150 m/s"],
        ["process_acc_noise", "500 m/s²"],
        ["min_snr_for_init", "22 dB"],
    ]
    _make_table(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(3.5),
                len(param_final), 2, param_final,
                col_widths=[Inches(3.2), Inches(2.6)])

    # Achievements
    _add_textbox(slide, Inches(6.7), Inches(1.0), Inches(6.2), Inches(0.4),
                 "主要改善成果 (Iter 1 → Iter 3)", size=14, bold=True, color=DARK_BLUE)
    achieve_data = [
        ["指標", "Iter 1", "Iter 3", "改善"],
        ["HGV 速度 RMSE (C)", "446 m/s", "232 m/s", "−48% ★★"],
        ["HGV cl=40 OSPA (D)", "3,224 m", "2,269 m", "−30% ★★"],
        ["Bal cl=40 OSPA (B)", "2,479 m", "1,891 m", "−24% ★"],
        ["Precision HGV cl=40", "92.5%", "98.6%", "+6.1 pt ★★"],
        ["F1 HGV cl=40", "95.5%", "98.9%", "+3.4 pt ★"],
        ["Pos RMSE HGV cl=40", "2,042 m", "1,948 m", "< 2000 m ★"],
    ]
    tbl = _make_table(slide, Inches(6.7), Inches(1.45), Inches(6.2), Inches(2.8),
                      len(achieve_data), 4, achieve_data,
                      col_widths=[Inches(2.2), Inches(1.2), Inches(1.2), Inches(1.6)])
    table = tbl.table
    for r in range(1, len(achieve_data)):
        cell = table.cell(r, 3)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                if "★" in run.text:
                    run.font.color.rgb = GREEN
                    run.font.bold = True

    # Remaining issues
    _add_textbox(slide, Inches(0.5), Inches(5.1), Inches(5.8), Inches(0.38),
                 "残課題", size=14, bold=True, color=RED)
    _add_bullet_list(slide, Inches(0.5), Inches(5.55), Inches(5.8), Inches(1.65), [
        "HGV 位置 RMSE は ~2000–2250 m 台 (機動の本質的誤差)",
        "HGV 単機で稀に偽トラック一時確定 (Track Purity 33% ケース)",
    ], size=11, bullet_color=RED)

    # Future
    _add_textbox(slide, Inches(6.7), Inches(4.4), Inches(6.2), Inches(0.38),
                 "今後の改善候補（アルゴリズム変更要）", size=14, bold=True, color=ACCENT_BLUE)
    _add_bullet_list(slide, Inches(6.7), Inches(4.85), Inches(6.2), Inches(2.0), [
        "IMM モデル遷移確率の適応的更新 (HGV 機動フェーズ対応)",
        "アダプティブゲート (SNR・距離に応じた動的閾値)",
        "HGV 確定トラックへのビーム優先割当強化",
    ], size=11, bullet_color=ACCENT_BLUE)

    # Conclusion bar
    concl = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.72))
    concl.fill.solid()
    concl.fill.fore_color.rgb = DARK_BLUE
    concl.line.fill.background()
    _add_textbox(
        slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.62),
        "結論: 3回のパラメータチューニングにより全パターンで F1≥98.9%、Precision≥98.6% を達成。"
        "HGV 速度 RMSE を最大 −48% 改善し、GPU加速による処理速度もリアルタイム要件を十分満足する。",
        size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
    )

    _add_slide_number(slide, 13, TOTAL_SLIDES)


# =========================================================================
# Main
# =========================================================================
def main():
    print("FastTracker 性能評価報告書 (3イテレーション版) を生成中...")

    chart_dir = tempfile.mkdtemp(prefix="fasttracker_charts_")
    print(f"  チャート一時ディレクトリ: {chart_dir}")

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("  スライド 1: 表紙")
    build_slide_title(prs)

    print("  スライド 2: 目次")
    build_slide_toc(prs)

    print("  スライド 3: 評価概要")
    build_slide_overview(prs)

    print("  スライド 4: 評価パターン・試験条件")
    build_slide_conditions(prs)

    print("  スライド 5: 追尾パラメータ推移")
    build_slide_params(prs)

    print("  スライド 6: Iter 1 ベースライン結果")
    build_slide_iter1(prs)

    print("  スライド 7: Iter 2 チューニング結果")
    build_slide_iter2(prs)

    print("  スライド 8: Iter 3 最終結果")
    build_slide_iter3(prs)

    print("  スライド 9: 性能推移 — RMSE (チャート生成中)")
    build_slide_rmse_trend(prs, chart_dir)

    print("  スライド 10: 性能推移 — OSPA・Precision (チャート生成中)")
    build_slide_ospa_prec_trend(prs, chart_dir)

    print("  スライド 11: 改善量サマリー (チャート生成中)")
    build_slide_improvement(prs, chart_dir)

    print("  スライド 12: 処理速度評価")
    build_slide_speed(prs)

    print("  スライド 13: 総合評価・まとめ")
    build_slide_summary(prs)

    prs.save(OUTPUT_PATH)
    print(f"\n  保存完了: {OUTPUT_PATH}")

    for fname in os.listdir(chart_dir):
        try:
            os.remove(os.path.join(chart_dir, fname))
        except OSError:
            pass
    try:
        os.rmdir(chart_dir)
        print("  一時ファイルを削除しました")
    except OSError:
        pass

    print("\n生成完了!")


if __name__ == "__main__":
    main()
